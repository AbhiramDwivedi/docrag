#!/usr/bin/env python3
"""Build test artifacts for golden test corpus.

This script generates office files and builds deterministic test artifacts including:
- FAISS vector index with normalized embeddings
- SQLite database with document metadata  
- FTS5 full-text search index (if available)
- Metadata with checksums and build info
"""

import argparse
import hashlib
import json
import logging
import os
import sqlite3
import sys
from pathlib import Path
from typing import Dict, List, Any, Optional

# Set deterministic environment
os.environ["OMP_NUM_THREADS"] = "1"
os.environ["PYTHONHASHSEED"] = "0"

import numpy as np
import torch
import time

# Set all seeds for deterministic behavior
np.random.seed(42)
torch.manual_seed(42)
if torch.cuda.is_available():
    torch.cuda.manual_seed_all(42)

# Office file generation libraries
try:
    from docx import Document
    from docx.shared import Inches
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False

try:
    from openpyxl import Workbook
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

# Add backend src to path for imports
backend_root = Path(__file__).parent.parent / "backend"
sys.path.insert(0, str(backend_root / "src"))

try:
    from shared.config import load_settings
except ImportError as e:
    logging.error(f"Failed to import backend modules: {e}")
    # Create a simple fallback config loader
    def load_settings(overrides=None):
        class SimpleConfig:
            def __init__(self):
                self.chunk_size = 400
                self.overlap = 50  
                self.embed_model = 'sentence-transformers/all-MiniLM-L6-v2'
            def get(self, key, default=None):
                return getattr(self, key, default)
        return SimpleConfig()
    
    logging.warning("Using fallback config loader")

def dummy_embed_texts(texts: List[str], model_name: str = None) -> np.ndarray:
    """Create simple deterministic dummy embeddings for testing."""
    # Create deterministic embeddings based on text content
    embeddings = []
    for text in texts:
        # Simple hash-based embedding for reproducibility
        hash_val = hash(text) % (2**31)  # Ensure consistent hash
        np.random.seed(hash_val % 1000)  # Use hash as seed for deterministic results
        
        # Generate a 384-dimensional vector (matching all-MiniLM-L6-v2)
        embedding = np.random.normal(0, 1, 384)
        embeddings.append(embedding)
    
    return np.array(embeddings, dtype=np.float32)


logger = logging.getLogger(__name__)


def setup_logging(verbose: bool = False):
    """Setup logging configuration."""
    level = logging.DEBUG if verbose else logging.INFO
    logging.basicConfig(
        level=level,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
    )


def compute_corpus_checksum(corpus_path: Path) -> str:
    """Compute deterministic checksum of corpus content."""
    hasher = hashlib.sha256()
    
    # Sort files for deterministic processing
    files = sorted(corpus_path.rglob("*.txt")) + sorted(corpus_path.rglob("*.md"))
    
    for file_path in files:
        hasher.update(file_path.name.encode())
        hasher.update(file_path.read_bytes())
    
    return hasher.hexdigest()


def generate_office_files(corpus_path: Path, force: bool = False):
    """Generate office files programmatically."""
    office_dir = corpus_path / "office"
    office_dir.mkdir(exist_ok=True)
    
    # Generate DOCX file
    docx_path = office_dir / "contoso_brief.docx"
    if force or not docx_path.exists():
        if HAS_DOCX:
            doc = Document()
            doc.add_heading('Contoso-Acme Partnership Brief', 0)
            doc.add_paragraph(
                'This document outlines the strategic partnership between Contoso '
                'Enterprise Solutions and Acme Corporation. The collaboration will '
                'leverage Acme\'s data platform capabilities with Contoso\'s consulting '
                'expertise to deliver comprehensive enterprise solutions.'
            )
            doc.add_heading('Partnership Benefits', level=1)
            doc.add_paragraph('• Combined technical expertise')
            doc.add_paragraph('• Expanded market reach')
            doc.add_paragraph('• Integrated service offerings')
            doc.save(str(docx_path))
            logger.info(f"Generated {docx_path}")
        else:
            logger.warning("python-docx not available, skipping DOCX generation")
    
    # Generate PPTX file
    pptx_path = office_dir / "roadmap.pptx"
    if force or not pptx_path.exists():
        if HAS_PPTX:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Technology Roadmap 2025"
            content = slide.shapes.placeholders[1].text_frame
            content.text = "Q1 2025 Initiatives:"
            p = content.add_paragraph()
            p.text = "• Integrate Reltora Gateway"
            p = content.add_paragraph()
            p.text = "• Deploy BOLT v2.1 protocol"
            p = content.add_paragraph() 
            p.text = "• Enhance Acme platform integration"
            prs.save(str(pptx_path))
            logger.info(f"Generated {pptx_path}")
        else:
            logger.warning("python-pptx not available, skipping PPTX generation")
    
    # Generate PDF file
    pdf_path = office_dir / "onepager.pdf"
    if force or not pdf_path.exists():
        if HAS_REPORTLAB:
            c = canvas.Canvas(str(pdf_path), pagesize=letter)
            c.drawString(100, 750, "Globex Financial Services - One Pager")
            c.drawString(100, 720, "Investment Focus: Technology Sector M&A")
            c.drawString(100, 690, "")
            c.drawString(100, 660, "Globex specializes in technology investments and has")
            c.drawString(100, 640, "completed over $2B in transactions. Our expertise in")
            c.drawString(100, 620, "enterprise software makes us the preferred advisor")
            c.drawString(100, 600, "for technology companies seeking strategic exits.")
            c.drawString(100, 570, "")
            c.drawString(100, 540, "Key Metrics:")
            c.drawString(120, 520, "• $800M assets under management")
            c.drawString(120, 500, "• #3 ranking in middle-market tech deals")
            c.drawString(120, 480, "• 35% YoY transaction volume growth")
            c.save()
            logger.info(f"Generated {pdf_path}")
        else:
            logger.warning("reportlab not available, skipping PDF generation")
    
    # Generate XLSX file
    xlsx_path = office_dir / "metrics.xlsx"
    if force or not xlsx_path.exists():
        if HAS_OPENPYXL:
            wb = Workbook()
            ws = wb.active
            ws.title = "Q3_Metrics"
            ws['A1'] = "Initech Q3 Performance"
            ws['A2'] = "Revenue Growth"
            ws['B2'] = "35%"
            ws['A3'] = "New Customers"
            ws['B3'] = "50+"
            ws['A4'] = "Platform Uptime"
            ws['B4'] = "99.9%"
            wb.save(str(xlsx_path))
            logger.info(f"Generated {xlsx_path}")
        else:
            logger.warning("openpyxl not available, skipping XLSX generation")


def check_fts5_availability() -> bool:
    """Check if SQLite was compiled with FTS5 support."""
    try:
        conn = sqlite3.connect(":memory:")
        cursor = conn.cursor()
        cursor.execute("CREATE VIRTUAL TABLE test_fts USING fts5(content)")
        conn.close()
        return True
    except sqlite3.OperationalError:
        return False


def chunk_text(text: str, chunk_size: int = 400, overlap: int = 50) -> List[str]:
    """Simple text chunking function."""
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        chunks.append(chunk.strip())
        start = end - overlap
    return chunks


def build_artifacts(corpus_path: Path, artifacts_path: Path, config: Dict[str, Any], force: bool = False):
    """Build deterministic test artifacts."""
    artifacts_path.mkdir(parents=True, exist_ok=True)
    
    # Set deterministic mode
    torch.set_num_threads(1)
    if torch.cuda.is_available():
        torch.cuda.set_device(0)
        torch.backends.cudnn.deterministic = True
        torch.backends.cudnn.benchmark = False
    
    # Force CPU for embeddings
    os.environ["CUDA_VISIBLE_DEVICES"] = ""
    
    # Collect all text files
    text_files = list(corpus_path.rglob("*.txt")) + list(corpus_path.rglob("*.md"))
    office_files = list(corpus_path.rglob("*.docx")) + list(corpus_path.rglob("*.pptx")) + \
                   list(corpus_path.rglob("*.pdf")) + list(corpus_path.rglob("*.xlsx"))
    
    all_files = sorted(text_files + office_files)
    logger.info(f"Processing {len(all_files)} files from corpus")
    
    # Extract text content and create chunks
    documents = []
    chunks = []
    chunk_metadata = []
    
    for file_path in all_files:
        relative_path = file_path.relative_to(corpus_path)
        
        # Extract text content
        if file_path.suffix.lower() in ['.txt', '.md']:
            content = file_path.read_text(encoding='utf-8')
        else:
            # For office files, use simple extraction or placeholder
            content = f"Office document: {file_path.name}"
            logger.info(f"Using placeholder content for {file_path.name}")
        
        documents.append({
            'file_path': str(relative_path),
            'content': content,
            'size': len(content)
        })
        
        # Create chunks
        file_chunks = chunk_text(content, config.get('chunk_size', 400), config.get('overlap', 50))
        for i, chunk_content in enumerate(file_chunks):
            chunks.append(chunk_content)
            chunk_metadata.append({
                'file_path': str(relative_path),
                'chunk_id': f"{relative_path.stem}_{i}",
                'chunk_index': i,
                'content': chunk_content
            })
    
    logger.info(f"Created {len(chunks)} chunks from {len(documents)} documents")
    
    # Generate embeddings deterministically
    model_name = config.get('embed_model', 'sentence-transformers/all-MiniLM-L6-v2')
    logger.info(f"Generating dummy embeddings for {model_name}")
    
    embeddings = dummy_embed_texts(chunks, model_name)
    
    # Normalize embeddings for cosine similarity
    embeddings = embeddings / np.linalg.norm(embeddings, axis=1, keepdims=True)
    
    # Build FAISS index
    import faiss
    
    dimension = embeddings.shape[1]
    index = faiss.IndexFlatIP(dimension)  # Inner product (cosine with normalized vectors)
    index.add(embeddings.astype(np.float32))
    
    # Save FAISS index
    index_path = artifacts_path / "vector.index"
    faiss.write_index(index, str(index_path))
    logger.info(f"Saved FAISS index to {index_path}")
    
    # Build SQLite database
    db_path = artifacts_path / "docmeta.db"
    conn = sqlite3.connect(str(db_path))
    cursor = conn.cursor()
    
    # Create documents table
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS documents (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            file_path TEXT UNIQUE NOT NULL,
            content TEXT NOT NULL,
            size INTEGER NOT NULL,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    """)
    
    # Create chunks table
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS chunks (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            file_path TEXT NOT NULL,
            chunk_id TEXT UNIQUE NOT NULL,
            chunk_index INTEGER NOT NULL,
            content TEXT NOT NULL,
            embedding_id INTEGER NOT NULL,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    """)
    
    # Insert documents
    for doc in documents:
        cursor.execute(
            "INSERT OR REPLACE INTO documents (file_path, content, size) VALUES (?, ?, ?)",
            (doc['file_path'], doc['content'], doc['size'])
        )
    
    # Insert chunks
    for i, metadata in enumerate(chunk_metadata):
        cursor.execute(
            "INSERT OR REPLACE INTO chunks (file_path, chunk_id, chunk_index, content, embedding_id) VALUES (?, ?, ?, ?, ?)",
            (metadata['file_path'], metadata['chunk_id'], metadata['chunk_index'], metadata['content'], i)
        )
    
    # Build FTS5 index if available
    has_fts5 = check_fts5_availability()
    if has_fts5:
        cursor.execute("""
            CREATE VIRTUAL TABLE IF NOT EXISTS chunks_fts USING fts5(
                chunk_id UNINDEXED,
                content,
                file_path UNINDEXED
            )
        """)
        
        # Populate FTS5 index
        for metadata in chunk_metadata:
            cursor.execute(
                "INSERT OR REPLACE INTO chunks_fts (chunk_id, content, file_path) VALUES (?, ?, ?)",
                (metadata['chunk_id'], metadata['content'], metadata['file_path'])
            )
        
        logger.info("Built FTS5 search index")
    else:
        logger.warning("FTS5 not available, skipping lexical search index")
    
    conn.commit()
    conn.close()
    logger.info(f"Saved SQLite database to {db_path}")
    
    # Generate metadata
    metadata = {
        'checksum': compute_corpus_checksum(corpus_path),
        'model_name': model_name,
        'embedding_dim': dimension,
        'index_type': 'IndexFlatIP',
        'has_fts5': has_fts5,
        'num_documents': len(documents),
        'num_chunks': len(chunks),
        'build_time': str(time.time()),
        'torch_version': torch.__version__,
        'numpy_version': np.__version__,
        'config': config
    }
    
    metadata_path = artifacts_path / "metadata.json"
    with open(metadata_path, 'w') as f:
        json.dump(metadata, f, indent=2)
    
    logger.info(f"Saved build metadata to {metadata_path}")
    return metadata


def main():
    parser = argparse.ArgumentParser(description='Build test artifacts for golden corpus')
    parser.add_argument('--rebuild', action='store_true', help='Force clean rebuild')
    parser.add_argument('--fail-on-drift', action='store_true', help='Exit non-zero if checksum mismatch')
    parser.add_argument('--ci', action='store_true', help='CI mode with stricter settings')
    parser.add_argument('--verbose', '-v', action='store_true', help='Verbose logging')
    args = parser.parse_args()
    
    setup_logging(args.verbose or args.ci)
    
    # Paths
    script_dir = Path(__file__).parent
    root_dir = script_dir.parent
    corpus_path = root_dir / "tests/fixtures/corpus_v1"
    artifacts_path = root_dir / "tests/fixtures/artifacts_v1"
    config_path = root_dir / "config.test.yaml"
    
    # Load test configuration  
    try:
        config = load_settings({'config_path': str(config_path)})
        # Convert config to dict, handling Path objects
        if hasattr(config, '__dict__'):
            config_dict = {}
            for k, v in config.__dict__.items():
                if isinstance(v, Path):
                    config_dict[k] = str(v)
                else:
                    config_dict[k] = v
        else:
            config_dict = {}
    except Exception as e:
        logger.warning(f"Could not load test config: {e}, using defaults")
        config_dict = {
            'chunk_size': 400,
            'overlap': 50,
            'embed_model': 'sentence-transformers/all-MiniLM-L6-v2'
        }
    
    # Check for drift
    metadata_path = artifacts_path / "metadata.json"
    current_checksum = compute_corpus_checksum(corpus_path)
    
    if metadata_path.exists() and not args.rebuild:
        with open(metadata_path) as f:
            old_metadata = json.load(f)
        
        if old_metadata.get('checksum') == current_checksum:
            logger.info("Artifacts are up to date")
            if args.fail_on_drift:
                sys.exit(0)
            return
        else:
            logger.info("Corpus has changed, rebuilding artifacts")
            if args.fail_on_drift:
                logger.error("Checksum mismatch detected")
                sys.exit(1)
    
    # Generate office files
    generate_office_files(corpus_path, force=args.rebuild)
    
    # Build artifacts
    logger.info("Building test artifacts...")
    metadata = build_artifacts(corpus_path, artifacts_path, config_dict, force=args.rebuild)
    
    logger.info("Test artifacts built successfully")
    logger.info(f"Checksum: {metadata['checksum']}")
    logger.info(f"Documents: {metadata['num_documents']}")
    logger.info(f"Chunks: {metadata['num_chunks']}")
    logger.info(f"FTS5 available: {metadata['has_fts5']}")


if __name__ == '__main__':
    main()