"""
PPTX extractor for PowerPoint presentations.

Features:
- Slide-by-slide text extraction
- Extracts text from all shapes
- Preserves slide structure as separate units
"""
from pathlib import Path
from typing import List

from .base import BaseExtractor, Unit

class PPTXExtractor(BaseExtractor):
    """PowerPoint presentation extractor using python-pptx."""
    
    @property
    def supported_extensions(self) -> List[str]:
        return [".pptx", ".ppt"]
    
    def extract(self, path: Path) -> List[Unit]:
        """Extract text from PPTX using python-pptx."""
        try:
            from pptx import Presentation
            
            print(f"📊 Processing PowerPoint presentation: {path.name}")
            
            prs = Presentation(str(path))
            units: List[Unit] = []
            
            for slide_num, slide in enumerate(prs.slides):
                text_parts: List[str] = []
                
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "text"):
                            text = getattr(shape, "text", "")
                            if text and str(text).strip():
                                text_parts.append(str(text))
                    except Exception:
                        continue  # Skip shapes that can't be processed
                
                slide_text = "\n".join(text_parts)
                if slide_text.strip():
                    units.append((f"slide_{slide_num + 1}", slide_text))
            
            print(f"   ✅ Extracted text from {len(units)} slides")
            return units
            
        except Exception as e:
            self._log_error(path, e)
            return []
