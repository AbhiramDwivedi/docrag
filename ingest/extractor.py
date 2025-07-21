"""File parsers for PDF, DOCX, PPTX, XLSX."""
from pathlib import Path
from typing import List, Tuple, Union
import os
import re

try:
    import pandas as pd
except ImportError:
    pd = None

Unit = Tuple[str, str]  # logical unit id, text

# Global flag for processing all sheets (set by command line)
_process_all_sheets = False

def set_all_sheets_mode(enabled: bool):
    """Set global flag for processing all sheets in Excel files."""
    global _process_all_sheets
    _process_all_sheets = enabled

def extract_text(path: Path) -> List[Unit]:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(path)
    if ext == ".docx":
        return _extract_docx(path)
    if ext in {".pptx", ".ppt"}:
        return _extract_pptx(path)
    if ext in {".xlsx", ".xls"}:
        return _extract_xlsx(path)
    if ext == ".txt":
        return _extract_txt(path)
    return []

# ---- private helpers with TODOs ----
def _extract_pdf(path: Path) -> List[Unit]:
    """Extract text from PDF using PyMuPDF (fitz)."""
    import fitz  # PyMuPDF
    try:
        doc = fitz.open(path)
        units = []
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text()
            if text.strip():
                units.append((f"page_{page_num + 1}", text))
        doc.close()
        return units
    except Exception as e:
        print(f"Error extracting PDF {path}: {e}")
        return []

def _extract_docx(path: Path) -> List[Unit]:
    """Extract text from DOCX using python-docx."""
    try:
        from docx import Document
        doc = Document(str(path))  # Convert Path to string
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        full_text = "\n".join(text_parts)
        return [("document", full_text)] if full_text.strip() else []
    except Exception as e:
        print(f"Error extracting DOCX {path}: {e}")
        return []

def _extract_pptx(path: Path) -> List[Unit]:
    """Extract text from PPTX using python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))  # Convert Path to string
        units = []
        for slide_num, slide in enumerate(prs.slides):
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
            slide_text = "\n".join(text_parts)
            if slide_text.strip():
                units.append((f"slide_{slide_num + 1}", slide_text))
        return units
    except Exception as e:
        print(f"Error extracting PPTX {path}: {e}")
        return []

def _extract_xlsx(path: Path) -> List[Unit]:
    """Extract text from XLSX by processing all sheets individually with smart prioritization."""
    try:
        import pandas as pd
        import re
        
        # Check file size first
        file_size_mb = path.stat().st_size / (1024 * 1024)
        if file_size_mb > 100:  # Increased limit for complete processing
            print(f"Skipping very large Excel file {path.name} ({file_size_mb:.1f}MB)")
            return []
        
        print(f"📊 Processing Excel file: {path.name} ({file_size_mb:.1f}MB)")
        
        # Get all sheet names without loading data
        with pd.ExcelFile(path, engine='openpyxl') as excel_file:
            all_sheets = excel_file.sheet_names
            
        print(f"   Found {len(all_sheets)} sheets: {', '.join(all_sheets[:5])}" + 
              (f" and {len(all_sheets)-5} more..." if len(all_sheets) > 5 else ""))
        
        # Prioritize sheets if more than 15 (unless all-sheets mode is enabled)
        sheets_to_process = all_sheets if _process_all_sheets else _prioritize_sheets(all_sheets)
        skipped_sheets = []
        
        if len(all_sheets) > 15 and not _process_all_sheets:
            skipped_sheets = [s for s in all_sheets if s not in sheets_to_process]
            print(f"   ⚠️  Limiting to 15 most relevant sheets. {len(skipped_sheets)} sheets will be skipped.")
            print(f"   📋 Processing: {', '.join(sheets_to_process)}")
        elif _process_all_sheets and len(all_sheets) > 15:
            print(f"   🔄 ALL-SHEETS MODE: Processing all {len(all_sheets)} sheets")
        
        units = []
        processed_count = 0
        empty_count = 0
        error_count = 0
        
        # Process each prioritized sheet
        for i, sheet_name in enumerate(sheets_to_process, 1):
            try:
                print(f"   📄 Sheet {i}/{len(sheets_to_process)}: '{sheet_name}'", end=" -> ")
                
                # Load sheet with generous limits for complete data
                df = pd.read_excel(
                    path, 
                    sheet_name=sheet_name, 
                    nrows=2000,  # Generous limit for complete data
                    engine='openpyxl'
                )
                
                # Check if sheet is empty or meaningless
                if _is_empty_sheet(df):
                    print("empty, skipped")
                    empty_count += 1
                    continue
                
                # Convert to meaningful text representation
                sheet_text = _convert_sheet_to_text(df, sheet_name)
                
                if sheet_text.strip():
                    units.append((f"sheet_{sheet_name}", sheet_text))
                    print(f"extracted {len(sheet_text)} chars")
                    processed_count += 1
                else:
                    print("no content, skipped")
                    empty_count += 1
                    
            except Exception as sheet_error:
                print(f"error: {sheet_error}")
                error_count += 1
                continue
        
        # Summary report
        print(f"   ✅ Excel processing complete: {processed_count} sheets processed, " +
              f"{empty_count} empty, {error_count} errors")
        
        if skipped_sheets:
            print(f"   💡 To process all sheets in '{path.name}', run:")
            print(f"      python -m ingest.ingest --file-type xlsx --target \"{path.name}\" --all-sheets")
        
        return units
        
    except Exception as e:
        print(f"Error extracting XLSX {path}: {e}")
        return []

def _prioritize_sheets(sheet_names: List[str]) -> List[str]:
    """Prioritize sheets based on meaningful names, returning top 15."""
    if len(sheet_names) <= 15:
        return sheet_names
    
    # Priority scoring based on sheet names
    priority_patterns = {
        'summary|overview|dashboard|main|primary|index': 10,
        'data|content|detail|information|info': 8,
        'report|analysis|results|findings': 7,
        'budget|financial|cost|revenue|sales': 6,
        'schedule|timeline|plan|roadmap': 5,
        'config|setting|parameter|metadata': 4,
        'template|example|sample': 2,
        'sheet\d+|temp|tmp|test|backup': 1
    }
    
    scored_sheets = []
    for sheet in sheet_names:
        score = 3  # Default score
        sheet_lower = sheet.lower()
        
        for pattern, points in priority_patterns.items():
            if re.search(pattern, sheet_lower):
                score = max(score, points)
                break
        
        scored_sheets.append((score, sheet))
    
    # Sort by score (highest first), then by original order
    scored_sheets.sort(key=lambda x: (-x[0], sheet_names.index(x[1])))
    
    return [sheet for score, sheet in scored_sheets[:15]]

def _is_empty_sheet(df) -> bool:
    """Check if a sheet is empty or contains only meaningless data."""
    if df.empty:
        return True
    
    # Count non-null, non-empty cells
    non_empty_cells = 0
    total_cells = df.shape[0] * df.shape[1]
    
    for col in df.columns:
        non_empty_cells += df[col].notna().sum()
        # Also check for non-whitespace string content
        if df[col].dtype == 'object':
            non_empty_cells += df[col].astype(str).str.strip().ne('').sum()
    
    # Consider empty if less than 1% of cells have content or very few cells total
    if total_cells > 0:
        content_ratio = non_empty_cells / total_cells
        return content_ratio < 0.01 or non_empty_cells < 5
    
    return True

def _convert_sheet_to_text(df, sheet_name: str) -> str:
    """Convert DataFrame to meaningful text representation preserving relationships."""
    if df.empty:
        return ""
    
    # Clean column names
    df.columns = [str(col).strip() for col in df.columns]
    
    # Start with sheet identification
    text_parts = [f"=== SHEET: {sheet_name} ==="]
    
    # Add column headers
    headers = " | ".join(df.columns)
    text_parts.append(f"COLUMNS: {headers}")
    text_parts.append("-" * min(80, len(headers)))
    
    # Process rows with relationship preservation
    for idx, row in df.iterrows():
        row_parts = []
        for col_name, value in row.items():
            if pd.notna(value) and str(value).strip():
                # Preserve column-value relationships
                clean_value = str(value).strip()
                if clean_value:
                    row_parts.append(f"{col_name}: {clean_value}")
        
        if row_parts:
            text_parts.append(" | ".join(row_parts))
    
    # Add summary statistics for numerical data
    numeric_cols = df.select_dtypes(include=['number']).columns
    if len(numeric_cols) > 0:
        text_parts.append("\n=== SUMMARY STATISTICS ===")
        for col in numeric_cols:
            if df[col].notna().sum() > 0:
                stats = {
                    'count': df[col].count(),
                    'mean': df[col].mean(),
                    'sum': df[col].sum(),
                    'min': df[col].min(),
                    'max': df[col].max()
                }
                stats_text = ", ".join([f"{k}: {v:.2f}" if isinstance(v, float) else f"{k}: {v}" 
                                      for k, v in stats.items()])
                text_parts.append(f"{col} -> {stats_text}")
    
    return "\n".join(text_parts)

def _extract_txt(path: Path) -> List[Unit]:
    """Extract text from plain text files."""
    try:
        with open(path, 'r', encoding='utf-8') as f:
            content = f.read()
        return [("document", content)] if content.strip() else []
    except Exception as e:
        print(f"Error extracting TXT {path}: {e}")
        return []
